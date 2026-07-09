from fastapi import FastAPI, HTTPException, Query
from pydantic import BaseModel
from pathlib import Path
from dotenv import load_dotenv
from typing import Optional, List
from sqlalchemy import create_engine, Column, Integer, String, DateTime, Boolean, Float, Text, ForeignKey, inspect, func, event, text
from sqlalchemy.orm import declarative_base, Session, relationship
import os
import datetime
import uuid
import hashlib
import json
import logging
import logging.handlers
import mimetypes

load_dotenv()

# ---------------------------------------------------------------------------
# Logging — file-based with rotation.
# Log directory is set via Orvyn_LOG_DIR env var (passed by Electron at spawn),
# falling back to python-backend/logs/ for standalone dev.
# ---------------------------------------------------------------------------
_log_dir = os.getenv("Orvyn_LOG_DIR", os.path.join(os.path.dirname(os.path.dirname(__file__)), "logs"))
os.makedirs(_log_dir, exist_ok=True)

_file_handler = logging.handlers.RotatingFileHandler(
    os.path.join(_log_dir, "python.log"),
    maxBytes=5 * 1024 * 1024,  # 5 MB
    backupCount=5,
    encoding="utf-8",
)
_file_handler.setFormatter(
    logging.Formatter("[%(asctime)s] [%(levelname)s] [%(name)s] %(message)s")
)

logging.basicConfig(level=logging.INFO, handlers=[_file_handler, logging.StreamHandler()])

logger = logging.getLogger("Orvyn")

app = FastAPI(title="Orvyn AI Engine")

# ---------------------------------------------------------------------------
# In-memory engine registry — set by POST /init-db, cleared on app restart.
# All DB-dependent routes must check this before operating.
# ---------------------------------------------------------------------------
active_engine = None

# Accepted theme values — validated before any DB write.
_VALID_THEMES = {"light", "dark"}

# Allowed file extensions for registration.
_ALLOWED_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".pptx", ".txt", ".csv", ".png", ".jpg", ".jpeg"}

# Max files per registration request.
_MAX_FILES_PER_REQUEST = 100

# Max characters stored in extracted_text column.
_MAX_EXTRACTED_TEXT_LENGTH = 3000

# ---------------------------------------------------------------------------
# SQLAlchemy ORM — schema defined and owned exclusively by this Python backend.
# No other layer may create, alter, or drop these tables.
# ---------------------------------------------------------------------------
Base = declarative_base()


class UserMeta(Base):
    __tablename__ = "user_meta"

    id = Column(Integer, primary_key=True, autoincrement=True)
    mongo_user_id = Column(String, nullable=False, unique=True)
    created_at = Column(DateTime, default=datetime.datetime.utcnow)


class Settings(Base):
    __tablename__ = "settings"

    key = Column(String, primary_key=True)
    value = Column(String, nullable=True)


# ---------------------------------------------------------------------------
# Smart DataRoom models
# ---------------------------------------------------------------------------

def _generate_uuid():
    return str(uuid.uuid4())


class DataRoom(Base):
    __tablename__ = "datarooms"

    id = Column(String, primary_key=True, default=_generate_uuid)
    name = Column(String, nullable=False)
    description = Column(Text, nullable=True)
    is_starred = Column(Boolean, default=False)
    created_by_ai = Column(Boolean, default=False)
    status = Column(String, default="active")  # active | archived
    created_at = Column(DateTime, default=datetime.datetime.utcnow)
    updated_at = Column(DateTime, default=datetime.datetime.utcnow, onupdate=datetime.datetime.utcnow)

    folders = relationship("Folder", back_populates="dataroom", cascade="all, delete-orphan")
    files = relationship("File", back_populates="dataroom", cascade="all, delete-orphan")


class Folder(Base):
    __tablename__ = "folders"

    id = Column(String, primary_key=True, default=_generate_uuid)
    dataroom_id = Column(String, ForeignKey("datarooms.id", ondelete="CASCADE"), nullable=False)
    name = Column(String, nullable=False)
    context = Column(Text, nullable=False)
    parent_id = Column(String, ForeignKey("folders.id"), nullable=True)
    display_order = Column(Integer, default=0)
    created_by_ai = Column(Boolean, default=False)
    created_at = Column(DateTime, default=datetime.datetime.utcnow)
    updated_at = Column(DateTime, default=datetime.datetime.utcnow, onupdate=datetime.datetime.utcnow)

    dataroom = relationship("DataRoom", back_populates="folders")
    children = relationship("Folder", backref="parent", remote_side=[id])
    files = relationship("File", back_populates="folder")
    classifications = relationship("Classification", back_populates="folder", cascade="all, delete-orphan")


class File(Base):
    __tablename__ = "files"

    id = Column(String, primary_key=True, default=_generate_uuid)
    dataroom_id = Column(String, ForeignKey("datarooms.id", ondelete="CASCADE"), nullable=False)
    folder_id = Column(String, ForeignKey("folders.id", ondelete="SET NULL"), nullable=True)
    original_name = Column(String, nullable=False)
    original_path = Column(Text, nullable=False)
    file_extension = Column(String, nullable=False)
    mime_type = Column(String, nullable=True)
    size_bytes = Column(Integer, nullable=True)
    checksum = Column(String, nullable=True)
    extracted_text = Column(Text, nullable=True)
    status = Column(String, default="registered")  # registered | processing | classified | error
    # Copilot columns
    ai_summary = Column(Text, nullable=True)
    embedding_status = Column(String, default="none")  # none | pending | processing | complete | failed
    content_checksum = Column(Text, nullable=True)
    embedding_model = Column(String, nullable=True)
    indexed_file_size = Column(Integer, nullable=True)
    indexed_file_mtime = Column(Float, nullable=True)
    created_at = Column(DateTime, default=datetime.datetime.utcnow)
    updated_at = Column(DateTime, default=datetime.datetime.utcnow, onupdate=datetime.datetime.utcnow)

    dataroom = relationship("DataRoom", back_populates="files")
    folder = relationship("Folder", back_populates="files")
    classifications = relationship("Classification", back_populates="file", cascade="all, delete-orphan")


class Classification(Base):
    __tablename__ = "classifications"

    id = Column(String, primary_key=True, default=_generate_uuid)
    file_id = Column(String, ForeignKey("files.id", ondelete="CASCADE"), nullable=False)
    folder_id = Column(String, ForeignKey("folders.id", ondelete="CASCADE"), nullable=False)
    confidence = Column(Float, nullable=True)
    reasoning = Column(Text, nullable=True)
    classified_at = Column(DateTime, default=datetime.datetime.utcnow)

    file = relationship("File", back_populates="classifications")
    folder = relationship("Folder", back_populates="classifications")


# ---------------------------------------------------------------------------
# Copilot models — new tables for the Copilot feature
# ---------------------------------------------------------------------------

class IndexingJob(Base):
    __tablename__ = "indexing_jobs"

    id = Column(String, primary_key=True, default=_generate_uuid)
    file_id = Column(String, ForeignKey("files.id", ondelete="CASCADE"), nullable=False)
    dataroom_id = Column(String, nullable=False)
    status = Column(String, default="pending")  # pending | processing | complete | failed
    attempts = Column(Integer, default=0)
    error_message = Column(Text, nullable=True)
    created_at = Column(DateTime, default=datetime.datetime.utcnow)
    updated_at = Column(DateTime, default=datetime.datetime.utcnow, onupdate=datetime.datetime.utcnow)


class FileChunk(Base):
    __tablename__ = "file_chunks"

    id = Column(String, primary_key=True, default=_generate_uuid)
    file_id = Column(String, ForeignKey("files.id", ondelete="CASCADE"), nullable=False)
    dataroom_id = Column(String, nullable=False)
    chunk_index = Column(Integer, nullable=False)
    chunk_text = Column(Text, nullable=False)


class ChatSession(Base):
    __tablename__ = "chat_sessions"

    id = Column(String, primary_key=True, default=_generate_uuid)
    scope_type = Column(String, nullable=False)
    scope_ids = Column(Text, nullable=True)
    scope_name = Column(Text, nullable=True)
    title = Column(Text, nullable=True)
    created_at = Column(DateTime, default=datetime.datetime.utcnow)
    updated_at = Column(DateTime, default=datetime.datetime.utcnow, onupdate=datetime.datetime.utcnow)


class ChatMessage(Base):
    __tablename__ = "chat_messages"

    id = Column(String, primary_key=True, default=_generate_uuid)
    session_id = Column(String, ForeignKey("chat_sessions.id", ondelete="CASCADE"), nullable=False)
    role = Column(String, nullable=False)
    content = Column(Text, nullable=False)
    sources = Column(Text, nullable=True)
    tool_calls = Column(Text, nullable=True)
    created_at = Column(DateTime, default=datetime.datetime.utcnow)


class FileEntity(Base):
    __tablename__ = "file_entities"

    id = Column(String, primary_key=True, default=_generate_uuid)
    file_id = Column(String, ForeignKey("files.id", ondelete="CASCADE"), nullable=False)
    dataroom_id = Column(String, nullable=False)
    entity_type = Column(String, nullable=False)
    entity_value = Column(String, nullable=False)
    context = Column(Text, nullable=True)


# ---------------------------------------------------------------------------
# Internal helpers
# ---------------------------------------------------------------------------

# Copilot context — set by /init-db, used by sync functions
_active_chroma_path = None
_active_user_id = None


def _get_user_id(session=None) -> str:
    """Return the active user's mongo_user_id for ChromaDB collection naming."""
    global _active_user_id
    if _active_user_id:
        return _active_user_id
    if session is not None:
        row = session.query(UserMeta).first()
        if row:
            _active_user_id = row.mongo_user_id
            return _active_user_id
    return None


def _get_chroma_path() -> str:
    """Return the ChromaDB storage path. Set via /init-db or copilot endpoints."""
    return _active_chroma_path


def _require_db():
    """
    Returns active_engine or raises 503 if /init-db has not been called yet.
    Call this at the start of every DB-dependent route.
    """
    if active_engine is None:
        raise HTTPException(
            status_code=503,
            detail="Database not initialised. Call POST /init-db first.",
        )
    return active_engine


# ---------------------------------------------------------------------------
# Text extraction helpers
# ---------------------------------------------------------------------------

def _compute_checksum(file_path: str) -> str:
    """Compute SHA-256 checksum of a file."""
    sha256 = hashlib.sha256()
    with open(file_path, "rb") as f:
        for chunk in iter(lambda: f.read(8192), b""):
            sha256.update(chunk)
    return sha256.hexdigest()


def _extract_text_pdf(file_path: str) -> str:
    """Extract text from PDF using PyMuPDF."""
    import fitz
    text_parts = []
    doc = fitz.open(file_path)
    for page in doc:
        text_parts.append(page.get_text())
    doc.close()
    return "\n".join(text_parts)


def _extract_text_docx(file_path: str) -> str:
    """Extract text from DOCX using python-docx."""
    from docx import Document
    doc = Document(file_path)
    return "\n".join(para.text for para in doc.paragraphs)


def _extract_text_xlsx(file_path: str) -> str:
    """Extract text from XLSX using openpyxl."""
    from openpyxl import load_workbook
    wb = load_workbook(file_path, read_only=True, data_only=True)
    text_parts = []
    for sheet_name in wb.sheetnames:
        text_parts.append(f"[Sheet: {sheet_name}]")
        ws = wb[sheet_name]
        for row in ws.iter_rows(values_only=True):
            cell_values = [str(cell) for cell in row if cell is not None]
            if cell_values:
                text_parts.append(" | ".join(cell_values))
    wb.close()
    return "\n".join(text_parts)


def _extract_text_pptx(file_path: str) -> str:
    """Extract text from PPTX using python-pptx."""
    from pptx import Presentation
    prs = Presentation(file_path)
    text_parts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        text_parts.append(f"[Slide {slide_num}]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        text_parts.append(para.text)
    return "\n".join(text_parts)


def _extract_text_plain(file_path: str) -> str:
    """Extract text from TXT/CSV with encoding fallback."""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    except UnicodeDecodeError:
        with open(file_path, "r", encoding="latin-1") as f:
            return f.read()


def _extract_text(file_path: str, extension: str, filename: str) -> str:
    """
    Extract text from a file based on its extension.
    Returns the extracted text or raises an exception on failure.
    """
    ext = extension.lower()

    if ext == ".pdf":
        return _extract_text_pdf(file_path)
    elif ext == ".docx":
        return _extract_text_docx(file_path)
    elif ext == ".xlsx":
        return _extract_text_xlsx(file_path)
    elif ext == ".pptx":
        return _extract_text_pptx(file_path)
    elif ext in (".txt", ".csv"):
        return _extract_text_plain(file_path)
    elif ext in (".png", ".jpg", ".jpeg"):
        return f"[Image: {filename}]"
    else:
        return ""


# ---------------------------------------------------------------------------
# Routes
# ---------------------------------------------------------------------------

@app.get("/health")
def health():
    result = {"status": "ok"}
    if active_engine is not None:
        inspector = inspect(active_engine)
        result["tables"] = inspector.get_table_names()
    return result


# ---- Database initialisation -----------------------------------------------

class InitDbRequest(BaseModel):
    database_path: str
    mongo_user_id: str


@app.post("/init-db")
def init_db(request: InitDbRequest):
    global active_engine, _active_chroma_path, _active_user_id

    db_path = request.database_path

    # Security: only absolute paths are accepted
    if not os.path.isabs(db_path):
        raise HTTPException(
            status_code=400,
            detail="database_path must be an absolute path",
        )

    # Security: reject any path containing '..' components
    path_obj = Path(db_path)
    if ".." in path_obj.parts:
        raise HTTPException(
            status_code=400,
            detail="Path traversal is not allowed in database_path",
        )

    # Resolve to normalise redundant separators without changing the target
    try:
        resolved = str(path_obj.resolve())
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid database_path")

    # Create parent directory if it does not exist
    parent_dir = os.path.dirname(resolved)
    try:
        os.makedirs(parent_dir, exist_ok=True)
    except OSError as exc:
        raise HTTPException(
            status_code=500,
            detail=f"Could not create database directory: {exc}",
        )

    # Create engine and initialise schema — Python owns all DDL
    engine = create_engine(f"sqlite:///{resolved}", echo=False)

    # SQLite requires an explicit pragma to enforce foreign key constraints.
    # Without this, ON DELETE CASCADE / SET NULL are silently ignored.
    @event.listens_for(engine, "connect")
    def _set_sqlite_pragma(dbapi_conn, _connection_record):
        cursor = dbapi_conn.cursor()
        cursor.execute("PRAGMA foreign_keys = ON")
        cursor.close()

    # Schema migration: if the files table exists with stale columns
    # (stored_name, local_path from old schema), drop files + classifications
    # so create_all rebuilds them with the current ORM definitions.
    with engine.connect() as conn:
        result = conn.execute(text(
            "SELECT name FROM sqlite_master WHERE type='table' AND name='files'"
        ))
        if result.fetchone():
            columns = [
                row[1] for row in conn.execute(text("PRAGMA table_info(files)"))
            ]
            if "original_path" not in columns:
                conn.execute(text("DROP TABLE IF EXISTS classifications"))
                conn.execute(text("DROP TABLE IF EXISTS files"))
                conn.commit()

    # Schema migration: add is_starred column to datarooms if missing
    with engine.connect() as conn:
        result = conn.execute(text(
            "SELECT name FROM sqlite_master WHERE type='table' AND name='datarooms'"
        ))
        if result.fetchone():
            columns = [
                row[1] for row in conn.execute(text("PRAGMA table_info(datarooms)"))
            ]
            if "is_starred" not in columns:
                conn.execute(text("ALTER TABLE datarooms ADD COLUMN is_starred BOOLEAN DEFAULT 0"))
                conn.commit()

    # Schema migration: add Copilot columns to files table if missing
    with engine.connect() as conn:
        result = conn.execute(text(
            "SELECT name FROM sqlite_master WHERE type='table' AND name='files'"
        ))
        if result.fetchone():
            columns = [
                row[1] for row in conn.execute(text("PRAGMA table_info(files)"))
            ]
            copilot_columns = {
                "ai_summary": "TEXT",
                "embedding_status": "TEXT DEFAULT 'none'",
                "content_checksum": "TEXT",
                "embedding_model": "TEXT",
                "indexed_file_size": "INTEGER",
                "indexed_file_mtime": "REAL",
            }
            for col_name, col_type in copilot_columns.items():
                if col_name not in columns:
                    conn.execute(text(f"ALTER TABLE files ADD COLUMN {col_name} {col_type}"))
            conn.commit()

    Base.metadata.create_all(engine)

    # Create FTS5 virtual table and sync triggers (content= mode)
    # These cannot be created via SQLAlchemy ORM — must use raw SQL.
    with engine.connect() as conn:
        # FTS5 virtual table
        conn.execute(text("""
            CREATE VIRTUAL TABLE IF NOT EXISTS file_chunks_fts USING fts5(
                chunk_text,
                file_id UNINDEXED,
                dataroom_id UNINDEXED,
                chunk_index UNINDEXED,
                content='file_chunks'
            )
        """))

        # AFTER INSERT trigger
        conn.execute(text("""
            CREATE TRIGGER IF NOT EXISTS file_chunks_ai AFTER INSERT ON file_chunks BEGIN
                INSERT INTO file_chunks_fts(rowid, chunk_text, file_id, dataroom_id, chunk_index)
                VALUES (new.rowid, new.chunk_text, new.file_id, new.dataroom_id, new.chunk_index);
            END
        """))

        # AFTER DELETE trigger
        conn.execute(text("""
            CREATE TRIGGER IF NOT EXISTS file_chunks_ad AFTER DELETE ON file_chunks BEGIN
                INSERT INTO file_chunks_fts(file_chunks_fts, rowid, chunk_text, file_id, dataroom_id, chunk_index)
                VALUES ('delete', old.rowid, old.chunk_text, old.file_id, old.dataroom_id, old.chunk_index);
            END
        """))

        # AFTER UPDATE trigger
        conn.execute(text("""
            CREATE TRIGGER IF NOT EXISTS file_chunks_au AFTER UPDATE ON file_chunks BEGIN
                INSERT INTO file_chunks_fts(file_chunks_fts, rowid, chunk_text, file_id, dataroom_id, chunk_index)
                VALUES ('delete', old.rowid, old.chunk_text, old.file_id, old.dataroom_id, old.chunk_index);
                INSERT INTO file_chunks_fts(rowid, chunk_text, file_id, dataroom_id, chunk_index)
                VALUES (new.rowid, new.chunk_text, new.file_id, new.dataroom_id, new.chunk_index);
            END
        """))

        # Performance index for recursive folder-tree CTE (Bug 7 perf fix)
        conn.execute(text(
            "CREATE INDEX IF NOT EXISTS idx_folders_parent_id ON folders(parent_id)"
        ))

        conn.commit()

    # Upsert user_meta row — idempotent, one row per mongo_user_id
    with Session(engine) as session:
        existing = (
            session.query(UserMeta)
            .filter_by(mongo_user_id=request.mongo_user_id)
            .first()
        )
        if not existing:
            session.add(UserMeta(mongo_user_id=request.mongo_user_id))
            session.commit()

    # Recover stale indexing jobs (crash recovery — Recommendation #6)
    try:
        with Session(engine) as session:
            from app.services.embedding_service import recover_stale_indexing_jobs
            recovered = recover_stale_indexing_jobs(session)
            if recovered:
                logger.info(f"init_db: recovered {recovered} stale indexing jobs")
    except Exception as e:
        logger.warning(f"init_db: indexing job recovery skipped (table may not exist yet): {e}")

    # Register engine only after schema + seed succeed
    active_engine = engine
    _active_user_id = request.mongo_user_id
    _active_chroma_path = os.path.join(os.path.dirname(resolved), "chroma")

    return {
        "status": "success",
        "message": "Database initialized",
        "path": resolved,
    }


# ---- Theme settings ---------------------------------------------------------

@app.get("/api/v1/settings/theme")
def get_theme():
    """
    Returns the stored theme for the active user.
    Defaults to "light" if no theme has been persisted yet.
    Requires /init-db to have been called first.
    """
    engine = _require_db()
    with Session(engine) as session:
        row = session.query(Settings).filter_by(key="theme").first()
    return {"theme": row.value if row else "light"}


class ThemeRequest(BaseModel):
    theme: str


@app.post("/api/v1/settings/theme")
def set_theme(request: ThemeRequest):
    """
    Persists the theme to the settings table.
    Only "light" and "dark" are accepted — all other values are rejected.
    Upserts the existing row so there is never more than one theme entry.
    Requires /init-db to have been called first.
    """
    if request.theme not in _VALID_THEMES:
        raise HTTPException(
            status_code=400,
            detail=f"Invalid theme value '{request.theme}'. Allowed: light, dark.",
        )

    engine = _require_db()
    with Session(engine) as session:
        row = session.query(Settings).filter_by(key="theme").first()
        if row:
            row.value = request.theme
        else:
            session.add(Settings(key="theme", value=request.theme))
        session.commit()

    return {"status": "success", "theme": request.theme}


# ---- DataRoom & Folder CRUD ------------------------------------------------

# -- Pydantic request models --

class CreateDataRoomRequest(BaseModel):
    name: str
    description: Optional[str] = None


class UpdateDataRoomRequest(BaseModel):
    name: Optional[str] = None
    description: Optional[str] = None
    is_starred: Optional[bool] = None


class CreateFolderRequest(BaseModel):
    name: str
    context: str
    parent_id: Optional[str] = None


class UpdateFolderRequest(BaseModel):
    name: Optional[str] = None
    context: Optional[str] = None
    parent_id: Optional[str] = "__unset__"  # distinguish between "not provided" and "set to null"


# -- Pydantic request models for Files --

class RegisterFilesRequest(BaseModel):
    dataroom_id: str
    file_paths: List[str]


class RelocateFileRequest(BaseModel):
    new_path: str


class MoveToFolderRequest(BaseModel):
    folder_id: Optional[str] = None
    dataroom_id: Optional[str] = None


class RenameFileRequest(BaseModel):
    new_name: str
    new_path: Optional[str] = None


# -- Serialisation helpers --

def _dt(val):
    """Convert a datetime to ISO-8601 string, or None."""
    return val.isoformat() if val else None


def _dataroom_dict(dr):
    return {
        "id": dr.id,
        "name": dr.name,
        "description": dr.description,
        "is_starred": bool(dr.is_starred),
        "created_by_ai": dr.created_by_ai,
        "status": dr.status,
        "created_at": _dt(dr.created_at),
        "updated_at": _dt(dr.updated_at),
    }


def _folder_dict(f):
    return {
        "id": f.id,
        "dataroom_id": f.dataroom_id,
        "name": f.name,
        "context": f.context,
        "parent_id": f.parent_id,
        "display_order": f.display_order,
        "created_by_ai": f.created_by_ai,
        "created_at": _dt(f.created_at),
        "updated_at": _dt(f.updated_at),
    }


def _file_dict(f):
    return {
        "id": f.id,
        "dataroom_id": f.dataroom_id,
        "folder_id": f.folder_id,
        "original_name": f.original_name,
        "original_path": f.original_path,
        "file_extension": f.file_extension,
        "mime_type": f.mime_type,
        "size_bytes": f.size_bytes,
        "checksum": f.checksum,
        "status": f.status,
        "embedding_status": f.embedding_status,
        "ai_summary": f.ai_summary,
        "created_at": _dt(f.created_at),
        "updated_at": _dt(f.updated_at),
    }


# -- DataRoom endpoints --

@app.post("/api/v1/datarooms", status_code=201)
def create_dataroom(request: CreateDataRoomRequest):
    if not request.name or not request.name.strip():
        raise HTTPException(status_code=400, detail="DataRoom name is required.")

    engine = _require_db()
    trimmed_name = request.name.strip()

    with Session(engine) as session:
        existing = session.query(DataRoom).filter(
            func.lower(DataRoom.name) == trimmed_name.lower()
        ).first()
        if existing:
            raise HTTPException(
                status_code=409,
                detail="DataRoom with this name already exists",
            )

        dr = DataRoom(name=trimmed_name, description=request.description)
        session.add(dr)
        session.commit()
        session.refresh(dr)
        return _dataroom_dict(dr)


@app.get("/api/v1/datarooms")
def list_datarooms():
    engine = _require_db()
    with Session(engine) as session:
        rows = session.query(DataRoom).order_by(DataRoom.created_at.desc()).all()

        results = []
        for dr in rows:
            folder_count = session.query(func.count(Folder.id)).filter(Folder.dataroom_id == dr.id).scalar()
            file_count = session.query(func.count(File.id)).filter(File.dataroom_id == dr.id).scalar()
            d = _dataroom_dict(dr)
            d["folder_count"] = folder_count
            d["file_count"] = file_count
            results.append(d)

        return results


@app.get("/api/v1/datarooms/{dataroom_id}")
def get_dataroom(dataroom_id: str):
    engine = _require_db()
    with Session(engine) as session:
        dr = session.query(DataRoom).filter_by(id=dataroom_id).first()
        if not dr:
            raise HTTPException(status_code=404, detail="DataRoom not found.")

        folders_raw = session.query(Folder).filter_by(dataroom_id=dr.id).order_by(Folder.display_order).all()
        folders = []
        for f in folders_raw:
            fd = _folder_dict(f)
            fd["file_count"] = session.query(func.count(File.id)).filter(File.folder_id == f.id).scalar()
            folders.append(fd)

        files = [
            _file_dict(f) for f in
            session.query(File).filter_by(dataroom_id=dr.id).order_by(File.created_at.desc()).all()
        ]

        result = _dataroom_dict(dr)
        result["folders"] = folders
        result["files"] = files
        return result


@app.put("/api/v1/datarooms/{dataroom_id}")
def update_dataroom(dataroom_id: str, request: UpdateDataRoomRequest):
    engine = _require_db()
    with Session(engine) as session:
        dr = session.query(DataRoom).filter_by(id=dataroom_id).first()
        if not dr:
            raise HTTPException(status_code=404, detail="DataRoom not found.")

        if request.name is not None:
            if not request.name.strip():
                raise HTTPException(status_code=400, detail="DataRoom name cannot be empty.")
            trimmed_name = request.name.strip()
            existing = session.query(DataRoom).filter(
                DataRoom.id != dataroom_id,
                func.lower(DataRoom.name) == trimmed_name.lower(),
            ).first()
            if existing:
                raise HTTPException(
                    status_code=409,
                    detail=f"A DataRoom with this name already exists",
                )
            dr.name = trimmed_name

        if request.description is not None:
            dr.description = request.description

        if request.is_starred is not None:
            dr.is_starred = request.is_starred

        dr.updated_at = datetime.datetime.utcnow()
        session.commit()
        session.refresh(dr)
        return _dataroom_dict(dr)


@app.delete("/api/v1/datarooms/{dataroom_id}")
def delete_dataroom(dataroom_id: str):
    engine = _require_db()
    with Session(engine) as session:
        dr = session.query(DataRoom).filter_by(id=dataroom_id).first()
        if not dr:
            raise HTTPException(status_code=404, detail="DataRoom not found.")

        # Sync: clean up all Copilot data before deleting DataRoom
        try:
            from app.services.embedding_service import sync_dataroom_deleted
            user_id = _get_user_id(session)
            chroma_path = _get_chroma_path()
            if user_id and chroma_path:
                sync_dataroom_deleted(dataroom_id, user_id, chroma_path, session)
        except Exception as e:
            logger.warning(f"delete_dataroom: sync_dataroom_deleted failed: {e}")

        session.delete(dr)
        session.commit()
        return {"success": True, "deleted_id": dataroom_id}


# -- Folder endpoints --

@app.post("/api/v1/datarooms/{dataroom_id}/folders", status_code=201)
def create_folder(dataroom_id: str, request: CreateFolderRequest):
    if not request.name or not request.name.strip():
        raise HTTPException(status_code=400, detail="Folder name is required.")
    if not request.context or not request.context.strip():
        raise HTTPException(status_code=400, detail="Folder context is required.")

    engine = _require_db()
    with Session(engine) as session:
        dr = session.query(DataRoom).filter_by(id=dataroom_id).first()
        if not dr:
            raise HTTPException(status_code=404, detail="DataRoom not found.")

        # Validate parent folder if provided
        if request.parent_id is not None:
            parent = session.query(Folder).filter_by(id=request.parent_id).first()
            if not parent:
                raise HTTPException(status_code=404, detail="Parent folder not found.")
            if parent.dataroom_id != dataroom_id:
                raise HTTPException(status_code=400, detail="Parent folder does not belong to this DataRoom.")

        # Check for duplicate folder name at the same parent level
        trimmed_name = request.name.strip()
        existing = session.query(Folder).filter(
            Folder.dataroom_id == dataroom_id,
            Folder.parent_id == request.parent_id,
            func.lower(Folder.name) == trimmed_name.lower(),
        ).first()
        if existing:
            raise HTTPException(
                status_code=409,
                detail="A folder with this name already exists here",
            )

        folder = Folder(
            dataroom_id=dataroom_id,
            name=trimmed_name,
            context=request.context.strip(),
            parent_id=request.parent_id,
        )
        session.add(folder)
        session.commit()
        session.refresh(folder)
        return _folder_dict(folder)


@app.get("/api/v1/datarooms/{dataroom_id}/folders")
def list_folders(dataroom_id: str):
    engine = _require_db()
    with Session(engine) as session:
        dr = session.query(DataRoom).filter_by(id=dataroom_id).first()
        if not dr:
            raise HTTPException(status_code=404, detail="DataRoom not found.")

        folders_raw = session.query(Folder).filter_by(dataroom_id=dataroom_id).order_by(Folder.display_order).all()
        results = []
        for f in folders_raw:
            fd = _folder_dict(f)
            fd["file_count"] = session.query(func.count(File.id)).filter(File.folder_id == f.id).scalar()
            results.append(fd)

        return results


@app.put("/api/v1/folders/{folder_id}")
def update_folder(folder_id: str, request: UpdateFolderRequest):
    engine = _require_db()
    with Session(engine) as session:
        folder = session.query(Folder).filter_by(id=folder_id).first()
        if not folder:
            raise HTTPException(status_code=404, detail="Folder not found.")

        if request.name is not None:
            if not request.name.strip():
                raise HTTPException(status_code=400, detail="Folder name cannot be empty.")
            trimmed_name = request.name.strip()
            existing = session.query(Folder).filter(
                Folder.id != folder_id,
                Folder.parent_id == folder.parent_id,
                Folder.dataroom_id == folder.dataroom_id,
                func.lower(Folder.name) == trimmed_name.lower(),
            ).first()
            if existing:
                raise HTTPException(
                    status_code=409,
                    detail=f"A folder with this name already exists here",
                )
            folder.name = trimmed_name

        if request.context is not None:
            if not request.context.strip():
                raise HTTPException(status_code=400, detail="Folder context cannot be empty.")
            folder.context = request.context.strip()

        # Handle parent_id change (folder move)
        if request.parent_id != "__unset__":
            if request.parent_id is not None:
                # Validate parent exists and belongs to same DataRoom
                parent = session.query(Folder).filter_by(id=request.parent_id).first()
                if not parent:
                    raise HTTPException(status_code=404, detail="Parent folder not found.")
                if parent.dataroom_id != folder.dataroom_id:
                    raise HTTPException(status_code=400, detail="Parent folder does not belong to the same DataRoom.")
                # Prevent circular reference — folder cannot be its own ancestor
                if request.parent_id == folder_id:
                    raise HTTPException(status_code=400, detail="A folder cannot be its own parent.")
            folder.parent_id = request.parent_id

        folder.updated_at = datetime.datetime.utcnow()
        session.commit()
        session.refresh(folder)
        return _folder_dict(folder)


@app.get("/api/v1/folders/{folder_id}/delete-preview")
def folder_delete_preview(folder_id: str):
    """Return counts of nested subfolders and files for the delete confirmation dialog."""
    engine = _require_db()
    with Session(engine) as session:
        folder = session.query(Folder).filter_by(id=folder_id).first()
        if not folder:
            raise HTTPException(status_code=404, detail="Folder not found.")

        # Recursive CTE to find all nested folder IDs
        all_folder_ids = _get_all_nested_folder_ids(session, folder_id)
        subfolder_count = len(all_folder_ids) - 1  # Exclude the target folder itself
        file_count = session.query(func.count(File.id)).filter(
            File.folder_id.in_(all_folder_ids)
        ).scalar()

        return {
            "subfolder_count": subfolder_count,
            "file_count": file_count,
        }


@app.delete("/api/v1/folders/{folder_id}")
def delete_folder(
    folder_id: str,
    file_action: str = Query(default="unassign", pattern="^(unassign|remove|delete_system)$"),
):
    """
    Delete a folder and all nested subfolders.

    file_action controls what happens to files inside:
      - unassign: files become unclassified (folder_id = NULL) — legacy default
      - remove: delete file records from DB only (disk untouched)
      - delete_system: delete files from disk AND from DB
    """
    engine = _require_db()
    with Session(engine) as session:
        folder = session.query(Folder).filter_by(id=folder_id).first()
        if not folder:
            raise HTTPException(status_code=404, detail="Folder not found.")

        # Recursive CTE to find ALL nested folder IDs
        all_folder_ids = _get_all_nested_folder_ids(session, folder_id)

        # Find all files in these folders
        affected_files = session.query(File).filter(
            File.folder_id.in_(all_folder_ids)
        ).all()

        # Sync: clean up Copilot data for all affected files
        dataroom_id = folder.dataroom_id
        try:
            from app.services.embedding_service import sync_folder_deleted
            user_id = _get_user_id(session)
            chroma_path = _get_chroma_path()
            if user_id and chroma_path and affected_files:
                all_file_ids = [f.id for f in affected_files]
                sync_folder_deleted(
                    folder_id, dataroom_id, all_file_ids,
                    user_id, chroma_path, session,
                )
        except Exception as e:
            logger.warning(f"delete_folder: sync_folder_deleted failed: {e}")

        files_deleted = 0
        files_removed = 0
        disk_errors = []

        if file_action == "delete_system":
            for f in affected_files:
                try:
                    if f.original_path and os.path.exists(f.original_path):
                        os.remove(f.original_path)
                    files_deleted += 1
                except OSError as e:
                    disk_errors.append(f"{f.original_name}: {str(e)}")
                session.delete(f)
        elif file_action == "remove":
            for f in affected_files:
                session.delete(f)
                files_removed += 1
        else:
            # unassign — legacy behavior: set folder_id to NULL
            for f in affected_files:
                f.folder_id = None

        # Delete all folders (children first by deleting all at once)
        session.query(Folder).filter(Folder.id.in_(all_folder_ids)).delete(
            synchronize_session="fetch"
        )
        session.commit()

        result = {
            "success": True,
            "folders_deleted": len(all_folder_ids),
        }
        if file_action == "delete_system":
            result["files_deleted"] = files_deleted
            if disk_errors:
                result["disk_errors"] = disk_errors
        elif file_action == "remove":
            result["files_removed"] = files_removed

        return result


def _get_all_nested_folder_ids(session, root_folder_id: str) -> list:
    """Use recursive CTE to find root folder + all descendant folder IDs."""
    cte = (
        session.query(Folder.id)
        .filter(Folder.id == root_folder_id)
        .cte(name="folder_tree", recursive=True)
    )
    cte = cte.union_all(
        session.query(Folder.id).filter(Folder.parent_id == cte.c.id)
    )
    rows = session.query(cte.c.id).all()
    return [r[0] for r in rows]


# ---- File endpoints ---------------------------------------------------------

@app.post("/api/v1/files/register", status_code=201)
def register_files(request: RegisterFilesRequest):
    """
    Register file paths and extract text content.
    Files stay at their original location — only the path is stored in SQLite.
    """
    engine = _require_db()

    # Validate file count limit
    if len(request.file_paths) > _MAX_FILES_PER_REQUEST:
        raise HTTPException(
            status_code=400,
            detail=f"Maximum {_MAX_FILES_PER_REQUEST} files per request. Received {len(request.file_paths)}.",
        )

    if not request.file_paths:
        raise HTTPException(status_code=400, detail="file_paths list cannot be empty.")

    with Session(engine) as session:
        # Verify DataRoom exists
        dr = session.query(DataRoom).filter_by(id=request.dataroom_id).first()
        if not dr:
            raise HTTPException(status_code=404, detail="DataRoom not found.")

        # Get all existing paths in this DataRoom for duplicate detection
        existing_paths = set(
            row[0] for row in
            session.query(File.original_path).filter_by(dataroom_id=request.dataroom_id).all()
        )

        registered = []
        rejected = []

        for file_path in request.file_paths:
            # Normalise path
            normalised = os.path.normpath(file_path)

            # Check: must be absolute
            if not os.path.isabs(normalised):
                rejected.append({"path": file_path, "reason": "Path must be absolute."})
                continue

            # Check: file extension
            _, ext = os.path.splitext(normalised)
            ext_lower = ext.lower()
            if ext_lower not in _ALLOWED_EXTENSIONS:
                rejected.append({"path": file_path, "reason": f"Extension '{ext_lower}' is not allowed. Allowed: {', '.join(sorted(_ALLOWED_EXTENSIONS))}"})
                continue

            # Check: file exists on disk
            if not os.path.exists(normalised):
                rejected.append({"path": file_path, "reason": "File does not exist at this path."})
                continue

            if not os.path.isfile(normalised):
                rejected.append({"path": file_path, "reason": "Path is not a file."})
                continue

            # Check: duplicate within this DataRoom
            if normalised in existing_paths:
                rejected.append({"path": file_path, "reason": "File already registered in this DataRoom."})
                continue

            # Read metadata
            original_name = os.path.basename(normalised)
            try:
                size_bytes = os.path.getsize(normalised)
            except OSError as exc:
                rejected.append({"path": file_path, "reason": f"Cannot read file size: {exc}"})
                continue

            mime_type, _ = mimetypes.guess_type(normalised)

            # Compute checksum
            try:
                checksum = _compute_checksum(normalised)
            except (OSError, PermissionError) as exc:
                rejected.append({"path": file_path, "reason": f"Cannot read file for checksum: {exc}"})
                continue

            # Extract text
            try:
                extracted = _extract_text(normalised, ext_lower, original_name)
            except Exception as exc:
                # File is corrupted or locked — register with error status
                file_record = File(
                    dataroom_id=request.dataroom_id,
                    original_name=original_name,
                    original_path=normalised,
                    file_extension=ext_lower,
                    mime_type=mime_type,
                    size_bytes=size_bytes,
                    checksum=checksum,
                    extracted_text=f"[Extraction error: {exc}]",
                    status="error",
                )
                session.add(file_record)
                session.flush()
                existing_paths.add(normalised)
                rejected.append({"path": file_path, "reason": f"Text extraction failed: {exc}"})
                continue

            # Truncate extracted text to max length
            if extracted and len(extracted) > _MAX_EXTRACTED_TEXT_LENGTH:
                extracted = extracted[:_MAX_EXTRACTED_TEXT_LENGTH]

            # Create file record
            file_record = File(
                dataroom_id=request.dataroom_id,
                original_name=original_name,
                original_path=normalised,
                file_extension=ext_lower,
                mime_type=mime_type,
                size_bytes=size_bytes,
                checksum=checksum,
                extracted_text=extracted,
                status="registered",
            )
            session.add(file_record)
            session.flush()
            session.refresh(file_record)
            existing_paths.add(normalised)

            registered.append(_file_dict(file_record))

        session.commit()

    return {
        "registered": registered,
        "rejected": rejected,
        "total_registered": len(registered),
        "total_rejected": len(rejected),
    }


@app.get("/api/v1/files/{file_id}")
def get_file(file_id: str):
    """Returns full file metadata including original_path."""
    engine = _require_db()
    with Session(engine) as session:
        file_record = session.query(File).filter_by(id=file_id).first()
        if not file_record:
            raise HTTPException(status_code=404, detail="File not found.")

        result = _file_dict(file_record)
        result["extracted_text"] = file_record.extracted_text
        return result


@app.post("/api/v1/files/{file_id}/check-exists")
def check_file_exists(file_id: str):
    """Checks if the file still exists at its original_path on disk."""
    engine = _require_db()
    with Session(engine) as session:
        file_record = session.query(File).filter_by(id=file_id).first()
        if not file_record:
            raise HTTPException(status_code=404, detail="File not found.")

        exists = os.path.exists(file_record.original_path) and os.path.isfile(file_record.original_path)
        return {"exists": exists, "path": file_record.original_path}


@app.put("/api/v1/files/{file_id}/relocate")
def relocate_file(file_id: str, request: RelocateFileRequest):
    """
    Update the stored path after a user has moved a file on disk.
    Validates the new path exists and checksum matches.
    """
    engine = _require_db()

    new_path = os.path.normpath(request.new_path)

    if not os.path.isabs(new_path):
        raise HTTPException(status_code=400, detail="new_path must be an absolute path.")

    if not os.path.exists(new_path) or not os.path.isfile(new_path):
        raise HTTPException(status_code=400, detail="File does not exist at the new path.")

    with Session(engine) as session:
        file_record = session.query(File).filter_by(id=file_id).first()
        if not file_record:
            raise HTTPException(status_code=404, detail="File not found.")

        # Verify checksum matches
        if file_record.checksum:
            try:
                new_checksum = _compute_checksum(new_path)
            except (OSError, PermissionError) as exc:
                raise HTTPException(status_code=400, detail=f"Cannot read file at new path: {exc}")

            if new_checksum != file_record.checksum:
                raise HTTPException(
                    status_code=400,
                    detail="Checksum mismatch. The file at the new path does not match the original file.",
                )

        file_record.original_path = new_path
        file_record.updated_at = datetime.datetime.utcnow()
        session.commit()
        session.refresh(file_record)

        # Sync: triple-check if file content has changed, re-index if so
        try:
            from app.services.embedding_service import has_file_changed, sync_file_content_changed
            user_id = _get_user_id(session)
            chroma_path = _get_chroma_path()
            if user_id and chroma_path and file_record.embedding_status == "complete":
                if has_file_changed(file_record, new_path):
                    logger.info(f"relocate_file: file {file_id} content changed, re-indexing")
                    sync_file_content_changed(
                        file_id, file_record.dataroom_id, user_id, chroma_path, session,
                    )
        except Exception as e:
            logger.warning(f"relocate_file: change detection failed: {e}")

        return _file_dict(file_record)


@app.put("/api/v1/files/{file_id}/move-to-folder")
def move_to_folder(file_id: str, request: MoveToFolderRequest):
    """Move file to a different virtual folder (null = unclassified)."""
    engine = _require_db()
    with Session(engine) as session:
        file_record = session.query(File).filter_by(id=file_id).first()
        if not file_record:
            raise HTTPException(status_code=404, detail="File not found.")

        if request.folder_id is not None:
            folder = session.query(Folder).filter_by(id=request.folder_id).first()
            if not folder:
                raise HTTPException(status_code=404, detail="Folder not found.")
            # If dataroom_id is provided, use it; otherwise check against file's current dataroom
            target_dataroom = request.dataroom_id or file_record.dataroom_id
            if folder.dataroom_id != target_dataroom:
                raise HTTPException(status_code=400, detail="Folder does not belong to the target DataRoom.")

        old_dataroom_id = file_record.dataroom_id
        old_folder_id = file_record.folder_id
        is_cross_dataroom = request.dataroom_id is not None and request.dataroom_id != old_dataroom_id

        file_record.folder_id = request.folder_id
        if request.dataroom_id is not None:
            file_record.dataroom_id = request.dataroom_id
        file_record.updated_at = datetime.datetime.utcnow()
        session.commit()
        session.refresh(file_record)

        # Sync: update ChromaDB metadata
        try:
            user_id = _get_user_id(session)
            chroma_path = _get_chroma_path()
            if user_id and chroma_path:
                if is_cross_dataroom:
                    from app.services.embedding_service import sync_file_moved_dataroom
                    sync_file_moved_dataroom(
                        file_id, old_dataroom_id, file_record.dataroom_id,
                        request.folder_id, user_id, chroma_path, session,
                    )
                else:
                    from app.services.embedding_service import sync_file_moved_folder
                    sync_file_moved_folder(
                        file_id, request.folder_id, user_id,
                        file_record.dataroom_id, chroma_path, session,
                    )
        except Exception as e:
            logger.warning(f"move_to_folder: sync failed: {e}")

        return _file_dict(file_record)


@app.delete("/api/v1/files/{file_id}")
def delete_file(file_id: str, delete_from_system: bool = Query(default=False)):
    """
    Remove file record from SQLite.
    If delete_from_system=true, also deletes the actual file from disk.
    """
    engine = _require_db()
    with Session(engine) as session:
        file_record = session.query(File).filter_by(id=file_id).first()
        if not file_record:
            raise HTTPException(status_code=404, detail="File not found.")

        file_dataroom_id = file_record.dataroom_id

        deleted_from_system = False
        if delete_from_system:
            try:
                if os.path.exists(file_record.original_path):
                    os.remove(file_record.original_path)
                    deleted_from_system = True
            except OSError as exc:
                raise HTTPException(
                    status_code=500,
                    detail=f"Failed to delete file from disk: {exc}",
                )

        # Sync: clean up Copilot data before deleting file record
        try:
            from app.services.embedding_service import sync_file_removed
            user_id = _get_user_id(session)
            chroma_path = _get_chroma_path()
            if user_id and chroma_path:
                sync_file_removed(file_id, user_id, file_dataroom_id, chroma_path, session)
        except Exception as e:
            logger.warning(f"delete_file: sync_file_removed failed: {e}")

        session.delete(file_record)
        session.commit()
        return {"success": True, "deleted_from_system": deleted_from_system}


@app.get("/api/v1/datarooms/{dataroom_id}/files")
def list_dataroom_files(
    dataroom_id: str,
    folder_id: Optional[str] = Query(default=None),
    include_subfolders: bool = Query(default=False),
    status: Optional[str] = Query(default=None),
):
    """
    List files in a DataRoom with optional filters.
    If folder_id is provided, filters to that folder.
    If include_subfolders is true, also includes files from nested child folders.
    """
    engine = _require_db()
    with Session(engine) as session:
        dr = session.query(DataRoom).filter_by(id=dataroom_id).first()
        if not dr:
            raise HTTPException(status_code=404, detail="DataRoom not found.")

        query = session.query(File).filter(File.dataroom_id == dataroom_id)

        if folder_id is not None:
            if include_subfolders:
                # Collect folder_id and all descendant folder IDs
                target_folder_ids = _collect_subfolder_ids(session, folder_id)
                query = query.filter(File.folder_id.in_(target_folder_ids))
            else:
                query = query.filter(File.folder_id == folder_id)

        if status is not None:
            query = query.filter(File.status == status)

        files = query.order_by(File.created_at.desc()).all()

        results = []
        for f in files:
            fd = _file_dict(f)
            if f.folder:
                fd["folder_name"] = f.folder.name
            else:
                fd["folder_name"] = None
            results.append(fd)

        return results


@app.put("/api/v1/files/{file_id}/rename")
def rename_file(file_id: str, request: RenameFileRequest):
    """
    Update display name in SQLite only — does NOT rename the actual file on disk.
    """
    engine = _require_db()

    if not request.new_name or not request.new_name.strip():
        raise HTTPException(status_code=400, detail="new_name cannot be empty.")

    with Session(engine) as session:
        file_record = session.query(File).filter_by(id=file_id).first()
        if not file_record:
            raise HTTPException(status_code=404, detail="File not found.")

        new_name = request.new_name.strip()
        file_record.original_name = new_name
        if request.new_path is not None:
            file_record.original_path = request.new_path
        file_record.updated_at = datetime.datetime.utcnow()
        session.commit()
        session.refresh(file_record)

        # Sync: update file_name in ChromaDB metadata
        try:
            from app.services.embedding_service import sync_file_renamed
            user_id = _get_user_id(session)
            chroma_path = _get_chroma_path()
            if user_id and chroma_path:
                sync_file_renamed(file_record.id, new_name, user_id, chroma_path)
        except Exception as e:
            logger.warning(f"rename_file: sync_file_renamed failed: {e}")

        return _file_dict(file_record)


# ---- OCR endpoints ----------------------------------------------------------
# Gemini Vision OCR is orchestrated by Electron:
#   Step 1: Python prepare-ocr → reads image bytes, base64-encodes
#   Step 2: Express /api/v1/ai/ocr → calls Gemini Vision, returns text
#   Step 3: Python apply-ocr → stores extracted text in SQLite

_IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg"}
_OCR_MAX_IMAGE_SIZE_BYTES = int(os.getenv("OCR_MAX_IMAGE_SIZE_MB", "10")) * 1024 * 1024


class PrepareOcrRequest(BaseModel):
    file_ids: List[str]


@app.post("/api/v1/files/prepare-ocr")
def prepare_ocr(request: PrepareOcrRequest):
    """
    Read image files from disk, base64-encode them for Gemini Vision OCR.
    Returns encoded image data for each file so Electron can forward to Express.
    """
    import base64

    engine = _require_db()

    if not request.file_ids:
        raise HTTPException(status_code=400, detail="file_ids list cannot be empty.")

    ocr_enabled = os.getenv("OCR_ENABLED", "true").lower() == "true"
    if not ocr_enabled:
        return {"files": [], "skipped": True, "reason": "OCR_ENABLED is false"}

    results = []

    with Session(engine) as session:
        for file_id in request.file_ids:
            file_record = session.query(File).filter_by(id=file_id).first()
            if not file_record:
                results.append({"file_id": file_id, "error": "File not found."})
                continue

            if file_record.file_extension not in _IMAGE_EXTENSIONS:
                results.append({"file_id": file_id, "error": "Not an image file."})
                continue

            file_path = file_record.original_path
            if not os.path.exists(file_path):
                results.append({"file_id": file_id, "error": "File not found on disk."})
                continue

            file_size = os.path.getsize(file_path)
            if file_size > _OCR_MAX_IMAGE_SIZE_BYTES:
                results.append({
                    "file_id": file_id,
                    "error": f"Image too large ({file_size // (1024*1024)}MB). Max {_OCR_MAX_IMAGE_SIZE_BYTES // (1024*1024)}MB.",
                })
                continue

            try:
                with open(file_path, "rb") as f:
                    image_bytes = f.read()
                image_base64 = base64.b64encode(image_bytes).decode("utf-8")
            except (OSError, PermissionError) as exc:
                results.append({"file_id": file_id, "error": f"Cannot read file: {exc}"})
                continue

            # Determine MIME type
            mime = file_record.mime_type or mimetypes.guess_type(file_path)[0] or "image/png"

            results.append({
                "file_id": file_id,
                "image_base64": image_base64,
                "mime_type": mime,
                "filename": file_record.original_name,
            })

    return {"files": results}


class ApplyOcrRequest(BaseModel):
    file_id: str
    extracted_text: str


@app.post("/api/v1/files/apply-ocr")
def apply_ocr(request: ApplyOcrRequest):
    """
    Store OCR-extracted text in the file's extracted_text column.
    If the file was already indexed, reset embedding_status so re-indexing picks it up.
    """
    engine = _require_db()

    with Session(engine) as session:
        file_record = session.query(File).filter_by(id=request.file_id).first()
        if not file_record:
            raise HTTPException(status_code=404, detail="File not found.")

        # Truncate to max length
        text = request.extracted_text.strip()
        if len(text) > _MAX_EXTRACTED_TEXT_LENGTH:
            text = text[:_MAX_EXTRACTED_TEXT_LENGTH]

        file_record.extracted_text = text
        file_record.updated_at = datetime.datetime.utcnow()

        # If file was already indexed (complete with 0 chunks from image skip),
        # reset so the indexing pipeline will re-process it with real text.
        if file_record.embedding_status in ("complete", "failed"):
            file_record.embedding_status = "none"
            file_record.content_checksum = None
            # Remove any existing indexing job so a fresh one can be created
            session.query(IndexingJob).filter_by(file_id=request.file_id).delete()
            logger.info(f"apply_ocr: reset embedding_status for file {request.file_id}")

        session.commit()

    logger.info(f"apply_ocr: stored {len(text)} chars for file {request.file_id}")
    return {"success": True, "file_id": request.file_id, "text_length": len(text)}


# ---- AI Data Preparation & Result Application endpoints ---------------------
# Gemini API calls have moved to the Express backend (holds the API key).
# Python now only prepares data (fingerprints, folder trees) and applies
# the AI results to the database. Electron orchestrates the full flow.

class PrepareClassifyRequest(BaseModel):
    dataroom_id: str
    file_ids: List[str]


class ApplyClassifyRequest(BaseModel):
    dataroom_id: str
    results: list  # Each: {file_id, folder_id, confidence, reasoning}


class PrepareGenerateRequest(BaseModel):
    file_ids: List[str]


class ApplyGenerateRequest(BaseModel):
    name: str
    description: Optional[str] = ""
    gemini_result: dict  # {folders: [...], assignments: [...]}
    file_ids: List[str]
    dataroom_id: Optional[str] = None


@app.post("/api/v1/ai/prepare-classify")
async def ai_prepare_classify(request: PrepareClassifyRequest):
    """
    Prepare fingerprints and folder tree for external AI classification.
    Does NOT call Gemini — returns data for Electron to forward to Express.
    """
    from app.services.classification_service import prepare_classify

    engine = _require_db()

    if not request.file_ids:
        raise HTTPException(status_code=400, detail="file_ids list cannot be empty.")

    if len(request.file_ids) > _MAX_FILES_PER_REQUEST:
        raise HTTPException(
            status_code=400,
            detail=f"Maximum {_MAX_FILES_PER_REQUEST} files per request. Received {len(request.file_ids)}.",
        )

    with Session(engine) as session:
        dr = session.query(DataRoom).filter_by(id=request.dataroom_id).first()
        if not dr:
            raise HTTPException(status_code=404, detail="DataRoom not found.")

    try:
        result = prepare_classify(engine, request.dataroom_id, request.file_ids)
        return result
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))


@app.post("/api/v1/ai/apply-classify")
async def ai_apply_classify(request: ApplyClassifyRequest):
    """
    Apply classification results (from Express/Gemini) to the database.
    Updates file folder assignments and creates Classification records.
    """
    from app.services.classification_service import apply_classify_results

    engine = _require_db()

    if not request.results:
        raise HTTPException(status_code=400, detail="results list cannot be empty.")

    with Session(engine) as session:
        dr = session.query(DataRoom).filter_by(id=request.dataroom_id).first()
        if not dr:
            raise HTTPException(status_code=404, detail="DataRoom not found.")

    try:
        result = apply_classify_results(engine, request.dataroom_id, request.results)

        # Create indexing jobs for ALL classified files (regardless of confidence)
        try:
            from app.services.embedding_service import create_indexing_job
            classified_file_ids = [
                r.get("file_id") for r in request.results
                if r.get("file_id")
            ]
            jobs_created = 0
            with Session(engine) as session:
                for file_id in classified_file_ids:
                    try:
                        create_indexing_job(file_id, request.dataroom_id, session)
                        jobs_created += 1
                    except Exception as e:
                        logger.warning(f"apply_classify: indexing job creation failed for {file_id}: {e}")
            if jobs_created:
                logger.info(f"apply_classify: created {jobs_created} indexing jobs")
        except Exception as e:
            logger.warning(f"apply_classify: indexing job setup failed: {e}")

        return result
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))


@app.post("/api/v1/ai/prepare-generate")
async def ai_prepare_generate(request: PrepareGenerateRequest):
    """
    Prepare file fingerprints for AI DataRoom generation.
    Does NOT call Gemini — returns data for Electron to forward to Express.
    """
    from app.services.classification_service import prepare_generate

    engine = _require_db()

    if not request.file_ids:
        raise HTTPException(status_code=400, detail="file_ids list cannot be empty.")

    if len(request.file_ids) > _MAX_FILES_PER_REQUEST:
        raise HTTPException(
            status_code=400,
            detail=f"Maximum {_MAX_FILES_PER_REQUEST} files per request. Received {len(request.file_ids)}.",
        )

    try:
        result = prepare_generate(engine, request.file_ids)
        return result
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))


@app.post("/api/v1/ai/apply-generate")
async def ai_apply_generate(request: ApplyGenerateRequest):
    """
    Apply AI-generated DataRoom structure (from Express/Gemini) to the database.
    Creates the DataRoom, folders, and assigns files.
    """
    from app.services.classification_service import apply_generate_results

    engine = _require_db()

    if not request.name or not request.name.strip():
        raise HTTPException(status_code=400, detail="name is required.")

    if not request.gemini_result:
        raise HTTPException(status_code=400, detail="gemini_result is required.")

    try:
        result = apply_generate_results(
            engine,
            request.name.strip(),
            request.description,
            request.gemini_result,
            request.file_ids,
            request.dataroom_id,
        )

        # Create indexing jobs for ALL assigned files (regardless of confidence)
        try:
            from app.services.embedding_service import create_indexing_job
            dataroom_id = result.get("dataroom", {}).get("id")
            if dataroom_id:
                assigned_file_ids = [
                    a.get("file_id") for a in request.gemini_result.get("assignments", [])
                    if a.get("file_id")
                ]
                jobs_created = 0
                with Session(engine) as session:
                    for file_id in assigned_file_ids:
                        try:
                            create_indexing_job(file_id, dataroom_id, session)
                            jobs_created += 1
                        except Exception as e:
                            logger.warning(f"apply_generate: indexing job creation failed for {file_id}: {e}")
                if jobs_created:
                    logger.info(f"apply_generate: created {jobs_created} indexing jobs")
        except Exception as e:
            logger.warning(f"apply_generate: indexing job setup failed: {e}")

        return result
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))


# ---------------------------------------------------------------------------
# Copilot endpoints — Embedding, Search, Chat, Indexing
# ---------------------------------------------------------------------------

# -- Pydantic request models for Copilot --

class PrepareIndexRequest(BaseModel):
    file_ids: List[str]
    dataroom_id: str
    chroma_path: Optional[str] = None

class ApplyIndexRequest(BaseModel):
    file_id: str
    dataroom_id: str
    chunks: list
    vectors: list
    embedding_model: str
    content_checksum: str
    file_size_bytes: Optional[int] = None
    file_mtime: Optional[float] = None
    user_id: str
    chroma_path: str
    preview_text: Optional[str] = None

class ApplyEntitiesRequest(BaseModel):
    file_id: str
    dataroom_id: str
    entities: dict  # {organizations: [], people: [], monetary_values: [], ...}

class ApplySummaryRequest(BaseModel):
    file_id: str
    summary: str

class CopilotSearchRequest(BaseModel):
    query_vector: list
    query_text: str
    scope_type: Optional[str] = "global"
    scope_ids: Optional[List[str]] = None
    session_id: Optional[str] = None
    scope_name: Optional[str] = None
    user_id: str
    chroma_path: str

class SaveMessageRequest(BaseModel):
    session_id: str
    user_message: str
    assistant_response: str
    sources: Optional[str] = None
    tool_calls: Optional[str] = None

class TriggerIndexingRequest(BaseModel):
    file_ids: List[str]
    dataroom_id: str

class RetryFailedRequest(BaseModel):
    dataroom_id: Optional[str] = None

class MarkFailedRequest(BaseModel):
    file_id: str
    error_message: Optional[str] = None

class CreateChatSessionRequest(BaseModel):
    scope_type: str
    scope_ids: Optional[str] = None
    scope_name: Optional[str] = None
    title: Optional[str] = None


@app.post("/api/v1/copilot/prepare-index")
def copilot_prepare_index(request: PrepareIndexRequest):
    """Chunk files, compute checksums, detect duplicates."""
    from app.services.embedding_service import prepare_index

    engine = _require_db()

    if not request.file_ids:
        raise HTTPException(status_code=400, detail="file_ids list cannot be empty.")

    with Session(engine) as session:
        result = prepare_index(request.file_ids, request.dataroom_id, session)
        return result


@app.post("/api/v1/copilot/apply-index")
def copilot_apply_index(request: ApplyIndexRequest):
    """Store vectors in ChromaDB + chunks in FTS5, update files table."""
    from app.services.embedding_service import apply_index

    engine = _require_db()

    with Session(engine) as session:
        result = apply_index(
            file_id=request.file_id,
            dataroom_id=request.dataroom_id,
            chunks=request.chunks,
            vectors=request.vectors,
            embedding_model=request.embedding_model,
            content_checksum=request.content_checksum,
            file_size_bytes=request.file_size_bytes,
            file_mtime=request.file_mtime,
            user_id=request.user_id,
            chroma_path=request.chroma_path,
            db_session=session,
            preview_text=request.preview_text,
        )
        return result


@app.post("/api/v1/copilot/apply-entities")
def copilot_apply_entities(request: ApplyEntitiesRequest):
    """Store extracted entities in file_entities table."""
    engine = _require_db()

    with Session(engine) as session:
        # Delete existing entities for this file
        session.execute(
            text("DELETE FROM file_entities WHERE file_id = :fid"),
            {"fid": request.file_id},
        )

        # Insert new entities
        count = 0
        for entity_type, values in request.entities.items():
            if not isinstance(values, list):
                continue
            for val in values:
                if isinstance(val, str):
                    entity_value = val
                    context = None
                elif isinstance(val, dict):
                    entity_value = val.get("value", str(val))
                    context = val.get("context")
                else:
                    entity_value = str(val)
                    context = None

                session.execute(
                    text("""
                        INSERT INTO file_entities (id, file_id, dataroom_id, entity_type, entity_value, context)
                        VALUES (:id, :fid, :did, :etype, :eval, :ctx)
                    """),
                    {
                        "id": str(uuid.uuid4()),
                        "fid": request.file_id,
                        "did": request.dataroom_id,
                        "etype": entity_type,
                        "eval": entity_value,
                        "ctx": context,
                    },
                )
                count += 1

        session.commit()
        return {"success": True, "entities_stored": count}


@app.post("/api/v1/copilot/apply-summary")
def copilot_apply_summary(request: ApplySummaryRequest):
    """Store AI-generated summary on file record."""
    engine = _require_db()

    with Session(engine) as session:
        file_record = session.query(File).filter_by(id=request.file_id).first()
        if not file_record:
            raise HTTPException(status_code=404, detail="File not found.")

        file_record.ai_summary = request.summary
        file_record.updated_at = datetime.datetime.utcnow()
        session.commit()
        return {"success": True, "file_id": request.file_id}


@app.post("/api/v1/copilot/search")
def copilot_search(request: CopilotSearchRequest):
    """Hybrid search with session management: delegates to prepare_chat_context.

    Returns { formatted_chunks, history, sources, session_id, session_title }
    as expected by the Electron orchestration layer.
    """
    from app.services.chat_service import prepare_chat_context

    engine = _require_db()

    with Session(engine) as session:
        try:
            return prepare_chat_context(
                message=request.query_text,
                query_vector=request.query_vector,
                session_id=request.session_id,
                scope_type=request.scope_type or "global",
                scope_ids=request.scope_ids,
                scope_name=request.scope_name,
                user_id=request.user_id,
                db_session=session,
                chroma_path=request.chroma_path,
            )
        except ValueError as e:
            raise HTTPException(status_code=400, detail=str(e))


@app.post("/api/v1/copilot/save-message")
def copilot_save_message(request: SaveMessageRequest):
    """Persist chat messages to SQLite."""
    engine = _require_db()

    with Session(engine) as session:
        # Insert user message
        session.execute(
            text("""
                INSERT INTO chat_messages (id, session_id, role, content)
                VALUES (:id, :sid, 'user', :content)
            """),
            {"id": str(uuid.uuid4()), "sid": request.session_id, "content": request.user_message},
        )
        # Insert assistant message
        session.execute(
            text("""
                INSERT INTO chat_messages (id, session_id, role, content, sources, tool_calls)
                VALUES (:id, :sid, 'assistant', :content, :sources, :tools)
            """),
            {
                "id": str(uuid.uuid4()),
                "sid": request.session_id,
                "content": request.assistant_response,
                "sources": request.sources,
                "tools": request.tool_calls,
            },
        )
        # Update session.updated_at
        session.execute(
            text("UPDATE chat_sessions SET updated_at = CURRENT_TIMESTAMP WHERE id = :sid"),
            {"sid": request.session_id},
        )
        session.commit()
        return {"success": True}


# -- Indexing endpoints --

@app.post("/api/v1/indexing/trigger")
def indexing_trigger(request: TriggerIndexingRequest):
    """Create indexing jobs for given file IDs."""
    from app.services.embedding_service import create_indexing_job

    engine = _require_db()

    if not request.file_ids:
        raise HTTPException(status_code=400, detail="file_ids list cannot be empty.")

    jobs_created = []
    with Session(engine) as session:
        for file_id in request.file_ids:
            # Check file exists
            file_exists = session.execute(
                text("SELECT id FROM files WHERE id = :fid"),
                {"fid": file_id},
            ).fetchone()
            if file_exists:
                job_id = create_indexing_job(file_id, request.dataroom_id, session)
                jobs_created.append({"file_id": file_id, "job_id": job_id})

    return {"success": True, "jobs_created": len(jobs_created), "jobs": jobs_created}


@app.get("/api/v1/indexing/status")
def indexing_status(dataroom_id: str = Query(default=None)):
    """Get indexing status counts per DataRoom (or all)."""
    engine = _require_db()

    with Session(engine) as session:
        scope_clause = ""
        params = {}
        if dataroom_id:
            scope_clause = "WHERE dataroom_id = :did"
            params["did"] = dataroom_id

        rows = session.execute(
            text(f"""
                SELECT status, COUNT(*) as cnt
                FROM indexing_jobs
                {scope_clause}
                GROUP BY status
            """),
            params,
        ).fetchall()

        counts = {"total": 0, "pending": 0, "processing": 0, "complete": 0, "failed": 0}
        for row in rows:
            status_val = row[0]
            cnt = row[1]
            if status_val in counts:
                counts[status_val] = cnt
            counts["total"] += cnt

        return counts


@app.post("/api/v1/indexing/retry-failed")
def indexing_retry_failed(request: RetryFailedRequest):
    """Reset failed jobs to pending for retry."""
    engine = _require_db()
    dataroom_id = request.dataroom_id

    with Session(engine) as session:
        scope_clause = ""
        params = {}
        if dataroom_id:
            scope_clause = "AND dataroom_id = :did"
            params["did"] = dataroom_id

        result = session.execute(
            text(f"""
                UPDATE indexing_jobs
                SET status = 'pending', error_message = NULL, updated_at = CURRENT_TIMESTAMP
                WHERE status = 'failed'
                {scope_clause}
            """),
            params,
        )
        # Also reset embedding_status on the corresponding files
        session.execute(
            text(f"""
                UPDATE files SET embedding_status = 'pending'
                WHERE id IN (
                    SELECT file_id FROM indexing_jobs
                    WHERE status = 'pending'
                    {scope_clause}
                ) AND embedding_status = 'failed'
            """),
            params,
        )
        session.commit()
        return {"success": True, "jobs_reset": result.rowcount}


@app.post("/api/v1/indexing/mark-failed")
def indexing_mark_failed(request: MarkFailedRequest):
    """Mark a specific file's indexing job as failed (called by Electron on pipeline error)."""
    engine = _require_db()

    with Session(engine) as session:
        result = session.execute(
            text("""
                UPDATE indexing_jobs
                SET status = 'failed',
                    error_message = :err,
                    attempts = attempts + 1,
                    updated_at = CURRENT_TIMESTAMP
                WHERE file_id = :fid AND status IN ('pending', 'processing')
            """),
            {"fid": request.file_id, "err": request.error_message or "Unknown error"},
        )

        if result.rowcount > 0:
            session.execute(
                text("UPDATE files SET embedding_status = 'failed' WHERE id = :fid"),
                {"fid": request.file_id},
            )

        session.commit()
        return {"success": True, "jobs_updated": result.rowcount}


@app.get("/api/v1/indexing/pending-files")
def indexing_pending_files(dataroom_id: str = Query(default=None)):
    """List file_ids with pending or processing indexing jobs.
    Called by Electron on startup for crash-recovery auto-resume."""
    engine = _require_db()

    with Session(engine) as session:
        scope_clause = ""
        params = {}
        if dataroom_id:
            scope_clause = "AND dataroom_id = :did"
            params["did"] = dataroom_id

        rows = session.execute(
            text(f"""
                SELECT file_id, dataroom_id
                FROM indexing_jobs
                WHERE status IN ('pending', 'processing')
                {scope_clause}
                ORDER BY created_at ASC
            """),
            params,
        ).fetchall()

        return {
            "files": [{"file_id": row[0], "dataroom_id": row[1]} for row in rows]
        }


# -- Chat session endpoints --

@app.post("/api/v1/chat/sessions", status_code=201)
def create_chat_session(request: CreateChatSessionRequest):
    """Create a new chat session."""
    engine = _require_db()

    with Session(engine) as session:
        chat = ChatSession(
            id=str(uuid.uuid4()),
            scope_type=request.scope_type,
            scope_ids=request.scope_ids,
            scope_name=request.scope_name,
            title=request.title,
        )
        session.add(chat)
        session.commit()
        session.refresh(chat)
        return {
            "id": chat.id,
            "scope_type": chat.scope_type,
            "scope_ids": chat.scope_ids,
            "scope_name": chat.scope_name,
            "title": chat.title,
            "created_at": _dt(chat.created_at),
            "updated_at": _dt(chat.updated_at),
        }


@app.get("/api/v1/chat/sessions")
def list_chat_sessions(
    scope_type: Optional[str] = Query(default=None),
    scope_id: Optional[str] = Query(default=None),
):
    """List chat sessions, optionally filtered by scope."""
    engine = _require_db()

    with Session(engine) as session:
        query = session.query(ChatSession)
        if scope_type:
            query = query.filter(ChatSession.scope_type == scope_type)
        if scope_id:
            # Use json_each for exact matching in scope_ids JSON array
            query = query.filter(
                ChatSession.id.in_(
                    session.query(ChatSession.id).filter(
                        text("EXISTS (SELECT 1 FROM json_each(chat_sessions.scope_ids) WHERE json_each.value = :sid)")
                    ).params(sid=scope_id)
                )
            )

        sessions_list = query.order_by(ChatSession.updated_at.desc()).all()
        return [
            {
                "id": s.id,
                "scope_type": s.scope_type,
                "scope_ids": s.scope_ids,
                "scope_name": s.scope_name,
                "title": s.title,
                "created_at": _dt(s.created_at),
                "updated_at": _dt(s.updated_at),
            }
            for s in sessions_list
        ]


@app.get("/api/v1/chat/sessions/{session_id}/messages")
def get_chat_messages(session_id: str):
    """Get all messages for a chat session."""
    engine = _require_db()

    with Session(engine) as session:
        chat = session.query(ChatSession).filter_by(id=session_id).first()
        if not chat:
            raise HTTPException(status_code=404, detail="Chat session not found.")

        messages = session.execute(
            text("""
                SELECT id, role, content, sources, tool_calls, created_at
                FROM chat_messages
                WHERE session_id = :sid
                ORDER BY created_at ASC
            """),
            {"sid": session_id},
        ).fetchall()

        return {
            "session_id": session_id,
            "messages": [
                {
                    "id": m[0],
                    "role": m[1],
                    "content": m[2],
                    "sources": json.loads(m[3]) if m[3] else [],
                    "tool_calls": json.loads(m[4]) if m[4] else [],
                    "created_at": m[5].isoformat() if m[5] else None,
                }
                for m in messages
            ],
        }


@app.delete("/api/v1/chat/sessions/{session_id}")
def delete_chat_session(session_id: str):
    """Delete a chat session and all its messages."""
    engine = _require_db()

    with Session(engine) as session:
        chat = session.query(ChatSession).filter_by(id=session_id).first()
        if not chat:
            raise HTTPException(status_code=404, detail="Chat session not found.")

        session.delete(chat)
        session.commit()
        return {"success": True, "deleted_id": session_id}


# ---------------------------------------------------------------------------
# Phase C2 — Copilot Tool Endpoints, Audit, Insights, Suggestions
# ---------------------------------------------------------------------------

# -- Pydantic request models for Phase C2 --

class ToolSearchRequest(BaseModel):
    query_vector: list
    query_text: str
    scope_type: Optional[str] = "global"
    scope_ids: Optional[List[str]] = None
    user_id: str
    chroma_path: str

class ToolGetFileContentRequest(BaseModel):
    file_id: str

class ToolListFilesRequest(BaseModel):
    dataroom_id: str
    folder_id: Optional[str] = None

class ToolGetEntitiesRequest(BaseModel):
    scope_type: str  # "file" | "dataroom"
    scope_id: str

class ToolFindSimilarRequest(BaseModel):
    file_id: str
    representative_chunk_vector: list
    user_id: str
    chroma_path: str
    max_results: Optional[int] = 5

class ToolPrepareCompareRequest(BaseModel):
    file_ids: List[str]

class ToolPrepareSummarizeRequest(BaseModel):
    dataroom_id: str

class ToolPrepareExtractRequest(BaseModel):
    query: str
    dataroom_id: str
    query_vector: list
    user_id: str
    chroma_path: str

class UpdateSessionTitleRequest(BaseModel):
    session_id: str
    title: str

class PrepareChatContextRequest(BaseModel):
    message: str
    query_vector: list
    session_id: Optional[str] = None
    scope_type: str = "global"
    scope_ids: Optional[List[str]] = None
    scope_name: Optional[str] = None
    user_id: str
    chroma_path: str


# -- Tool endpoints (called by Electron when Gemini requests a tool) --

@app.post("/api/v1/copilot/tool/search")
def copilot_tool_search(request: ToolSearchRequest):
    """Tool: search_documents — hybrid search for document content."""
    from app.services.copilot_tools import tool_search_documents

    engine = _require_db()
    with Session(engine) as session:
        return tool_search_documents(
            query_vector=request.query_vector,
            query_text=request.query_text,
            scope_type=request.scope_type,
            scope_ids=request.scope_ids,
            user_id=request.user_id,
            db_session=session,
            chroma_path=request.chroma_path,
        )


@app.post("/api/v1/copilot/tool/get-file-content")
def copilot_tool_get_file_content(request: ToolGetFileContentRequest):
    """Tool: get_file_content — fetch extracted text of a specific file."""
    from app.services.copilot_tools import tool_get_file_content

    engine = _require_db()
    with Session(engine) as session:
        return tool_get_file_content(request.file_id, session)


@app.post("/api/v1/copilot/tool/list-files")
def copilot_tool_list_files(request: ToolListFilesRequest):
    """Tool: list_files_in_dataroom — list files with metadata."""
    from app.services.copilot_tools import tool_list_files

    engine = _require_db()
    with Session(engine) as session:
        return tool_list_files(request.dataroom_id, request.folder_id, session)


@app.post("/api/v1/copilot/tool/get-entities")
def copilot_tool_get_entities(request: ToolGetEntitiesRequest):
    """Tool: get_entities — query extracted entities grouped by type."""
    from app.services.copilot_tools import tool_get_entities

    engine = _require_db()
    with Session(engine) as session:
        return tool_get_entities(request.scope_type, request.scope_id, session)


@app.post("/api/v1/copilot/tool/find-similar")
def copilot_tool_find_similar(request: ToolFindSimilarRequest):
    """Tool: find_similar_documents — find similar docs across DataRooms."""
    from app.services.copilot_tools import tool_find_similar

    return tool_find_similar(
        file_id=request.file_id,
        representative_chunk_vector=request.representative_chunk_vector,
        user_id=request.user_id,
        chroma_path=request.chroma_path,
        max_results=request.max_results or 5,
    )


@app.post("/api/v1/copilot/tool/prepare-compare")
def copilot_tool_prepare_compare(request: ToolPrepareCompareRequest):
    """Tool: compare_documents — prepare file data for Gemini comparison."""
    from app.services.copilot_tools import prepare_compare_data

    engine = _require_db()
    with Session(engine) as session:
        return prepare_compare_data(request.file_ids, session)


@app.post("/api/v1/copilot/tool/prepare-summarize")
def copilot_tool_prepare_summarize(request: ToolPrepareSummarizeRequest):
    """Tool: summarize_dataroom — prepare DataRoom data for Gemini summary."""
    from app.services.copilot_tools import prepare_summarize_data

    engine = _require_db()
    with Session(engine) as session:
        return prepare_summarize_data(request.dataroom_id, session)


@app.post("/api/v1/copilot/tool/prepare-extract")
def copilot_tool_prepare_extract(request: ToolPrepareExtractRequest):
    """Tool: extract_data_point — search for specific data for Gemini extraction."""
    from app.services.copilot_tools import prepare_extract_data

    engine = _require_db()
    with Session(engine) as session:
        return prepare_extract_data(
            query=request.query,
            dataroom_id=request.dataroom_id,
            query_vector=request.query_vector,
            user_id=request.user_id,
            chroma_path=request.chroma_path,
            db_session=session,
        )


# -- Document comparison --

class PrepareCompareRequest(BaseModel):
    file_ids: List[str]


@app.post("/api/v1/copilot/prepare-compare")
def copilot_prepare_compare(request: PrepareCompareRequest):
    """
    Prepare structured document content for Gemini comparison.
    Returns { files: [{ file_id, file_name, file_type, content }], file_count }
    where content is capped at 3000 chars per file.
    """
    from app.services.copilot_tools import prepare_compare_data

    engine = _require_db()
    with Session(engine) as session:
        return prepare_compare_data(request.file_ids, session)


# -- Stale content detection --

class CheckFileChangedRequest(BaseModel):
    file_id: str


@app.post("/api/v1/copilot/check-file-changed")
def copilot_check_file_changed(request: CheckFileChangedRequest):
    """
    Check whether a file's on-disk content has changed since it was last indexed.
    Returns { "changed": true/false }.
    """
    from app.services.embedding_service import has_file_changed

    engine = _require_db()
    with Session(engine) as session:
        file_record = session.query(File).filter_by(id=request.file_id).first()
        if not file_record:
            raise HTTPException(status_code=404, detail="File not found.")

        changed = has_file_changed(file_record, file_record.original_path)
        return {"changed": changed}


# -- Chat context preparation endpoint (Phase C2) --

@app.post("/api/v1/copilot/prepare-chat")
def copilot_prepare_chat(request: PrepareChatContextRequest):
    """Prepare chat context: session, hybrid search, history, formatted chunks."""
    from app.services.chat_service import prepare_chat_context

    engine = _require_db()
    with Session(engine) as session:
        try:
            return prepare_chat_context(
                message=request.message,
                query_vector=request.query_vector,
                session_id=request.session_id,
                scope_type=request.scope_type,
                scope_ids=request.scope_ids,
                scope_name=request.scope_name,
                user_id=request.user_id,
                db_session=session,
                chroma_path=request.chroma_path,
            )
        except ValueError as e:
            raise HTTPException(status_code=400, detail=str(e))


@app.post("/api/v1/copilot/update-session-title")
def copilot_update_session_title(request: UpdateSessionTitleRequest):
    """Update the title of a chat session."""
    from app.services.chat_service import update_session_title

    engine = _require_db()
    with Session(engine) as session:
        update_session_title(request.session_id, request.title, session)
        return {"success": True, "session_id": request.session_id}


# ---------------------------------------------------------------------------
# Subfolder traversal helper
# ---------------------------------------------------------------------------

def _collect_subfolder_ids(session, folder_id: str) -> List[str]:
    """
    BFS traversal to collect a folder and all its descendant folder IDs.
    """
    result = [folder_id]
    queue = [folder_id]
    while queue:
        current_id = queue.pop(0)
        children = session.query(Folder.id).filter_by(parent_id=current_id).all()
        for (child_id,) in children:
            result.append(child_id)
            queue.append(child_id)
    return result
